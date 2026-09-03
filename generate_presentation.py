import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Installing python-pptx...")
    os.system(f"{sys.executable} -m pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

def create_creative_deck():
    prs = Presentation()
    # 16:9 widescreen layout (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Curated Palette
    BG_DARK = RGBColor(11, 17, 32)       # #0B1120
    CARD_BG = RGBColor(15, 23, 42)       # #0F172A
    CARD_BG_ALT = RGBColor(19, 31, 56)   # #131F38
    CARD_BORDER = RGBColor(30, 41, 59)   # #1E293B
    BORDER_GLOW = RGBColor(38, 70, 110)  # #26466E
    
    CYAN = RGBColor(6, 182, 212)         # #06B6D4
    BLUE = RGBColor(59, 130, 246)        # #3B82F6
    PURPLE = RGBColor(168, 85, 247)      # #A855F7
    RED = RGBColor(239, 68, 68)          # #EF4444
    AMBER = RGBColor(245, 158, 11)       # #F59E0B
    EMERALD = RGBColor(16, 185, 129)     # #10B981
    
    TEXT_WHITE = RGBColor(248, 250, 252) # #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8
    TEXT_DIM = RGBColor(100, 116, 139)   # #64748B

    # Image Paths
    IMG_DIR = r"C:\Users\tella\.gemini\antigravity-ide\brain\6afea106-ae31-40ff-b727-93b6922a5fbc"
    IMG_GRAPH = os.path.join(IMG_DIR, "fincrime_hero_network_1788422771759.jpg")
    IMG_BIO = os.path.join(IMG_DIR, "biometrics_telemetry_hud_1788422797683.jpg")
    IMG_WAR_ROOM = os.path.join(IMG_DIR, "fiu_war_room_1788422822847.jpg")

    def add_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()

    def add_header(slide, title_text, category_text="SALAAR BANK • NEXT-GEN FINCRIME OS"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_c = cat_box.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = CYAN

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_t = title_box.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    # =========================================================================
    # SLIDE 1: Hero Cover Slide (Split Screen with Embedded Visual)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    add_background(s1)

    # Left: Typography & Value Prop
    add_card(s1, Inches(0.8), Inches(0.8), Inches(5.8), Inches(5.9), bg_color=CARD_BG, border_color=BORDER_GLOW)
    
    # Accent badge
    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(1.1), Inches(3.2), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(8, 47, 73)
    badge.line.color.rgb = CYAN
    badge.line.width = Pt(1)
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "🏛️ INSTITUTIONAL DEFENSE SUITE"
    p_b.font.size = Pt(10)
    p_b.font.bold = True
    p_b.font.color.rgb = CYAN

    tbox = s1.shapes.add_textbox(Inches(1.1), Inches(1.65), Inches(5.2), Inches(4.7))
    tf = tbox.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = "SALAAR FinCrime OS"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(4)

    p1_sub = tf.add_paragraph()
    p1_sub.text = "Next-Gen Multi-Entity Financial Crime & Autonomous Defense Operating System"
    p1_sub.font.size = Pt(16)
    p1_sub.font.bold = True
    p1_sub.font.color.rgb = CYAN
    p1_sub.space_after = Pt(14)

    p2 = tf.add_paragraph()
    p2.text = "Eliminating legacy heuristics with sub-15ms multi-hop topology analysis, pre-transaction behavioral biometrics, and autonomous AI regulatory filing."
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_after = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "⚡ Stack: Next.js 16 • Supabase PostgreSQL • Cytoscape • BioCatch Engine • Multi-Agent AI"
    p3.font.size = Pt(11)
    p3.font.bold = True
    p3.font.color.rgb = EMERALD

    # Right: Embedded Hero Graphic Card
    if os.path.exists(IMG_GRAPH):
        add_card(s1, Inches(6.8), Inches(0.8), Inches(5.7), Inches(5.9), bg_color=CARD_BG, border_color=BORDER_GLOW)
        s1.shapes.add_picture(IMG_GRAPH, Inches(6.9), Inches(0.9), width=Inches(5.5), height=Inches(5.7))

    # =========================================================================
    # SLIDE 2: The Core Problem & Industry Crisis
    # =========================================================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    add_background(s2)
    add_header(s2, "The $5.8B Annual Crisis: Why Isolated Rules Fail Modern Banking")

    problems = [
        ("💸 Authorised Push Payment (APP) Scams", "Victims are socially engineered over phone or WhatsApp to directly copy-paste scammer accounts. Post-login rules have zero visibility into user coercion.", RED),
        ("🕸️ Rapid Drain Money Mule Networks", "Organized syndicates funnel funds through 5+ accounts and drain >=85% within <15 mins, far outpacing legacy T+1 batch compliance jobs.", AMBER),
        ("📱 Device Farms & Synthetic Personas", "Fraud rings operate 20+ accounts from identical browser canvas hardware fingerprints, rotating residential proxies to bypass naive IP checks.", PURPLE),
        ("⚖️ Devastating Regulatory Penalties", "Delayed FinCEN & RBI Suspicious Activity Reports (SAR/STR) result in massive compliance fines and operational license suspensions.", BLUE),
    ]

    card_w = Inches(5.65)
    card_h = Inches(2.35)
    for i, (head, desc, accent) in enumerate(problems):
        row = i // 2
        col = i % 2
        c_left = Inches(0.8) + col * (card_w + Inches(0.4))
        c_top = Inches(1.8) + row * (card_h + Inches(0.35))

        add_card(s2, c_left, c_top, card_w, card_h)

        # Indicator tag
        tag = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, c_left, c_top, Inches(0.1), card_h)
        tag.fill.solid()
        tag.fill.fore_color.rgb = accent
        tag.line.fill.background()

        box = s2.shapes.add_textbox(c_left + Inches(0.25), c_top + Inches(0.2), card_w - Inches(0.5), card_h - Inches(0.4))
        btf = box.text_frame
        btf.word_wrap = True

        hp = btf.paragraphs[0]
        hp.text = head
        hp.font.size = Pt(16)
        hp.font.bold = True
        hp.font.color.rgb = accent
        hp.space_after = Pt(8)

        dp = btf.add_paragraph()
        dp.text = desc
        dp.font.size = Pt(13)
        dp.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 3: System Architecture & The 4 Pillars
    # =========================================================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    add_background(s3)
    add_header(s3, "Multi-Vector Defense OS: 4 Interconnected Intelligence Pillars")

    pillars = [
        ("PILLAR 1", "Knowledge Graph", "Heterogeneous Multi-Entity", "Correlates Accounts, Users, Devices & Subnets.\n• Tarjan's Wash Loop Engine\n• Rapid Drain Mule Intercept\n• Capital Contagion Radius", CYAN),
        ("PILLAR 2", "Biometrics", "Pre-Tx Telemetry", "Evaluates risk while typing.\n• Keystroke Cadence Entropy\n• Clipboard Paste APP Alert\n• Haversine Impossible Travel", EMERALD),
        ("PILLAR 3", "Sanctions & PEP", "Watchlist Screener", "Compliance-grade checks.\n• Jaro-Winkler Phonetic Match\n• OFAC SDN & PEP Lists\n• ISO 20022 pacs.008 Parser", PURPLE),
        ("PILLAR 4", "Autonomous FIU", "AI Multi-Agent", "Autonomous containment.\n• 1-Click Cluster Quarantine\n• Auto-SAR Legal Filing\n• BSA 31 U.S.C. 5318(g)", BLUE)
    ]

    p_w = Inches(2.7)
    p_h = Inches(5.1)
    for i, (badge, title, subtitle, desc, acc) in enumerate(pillars):
        c_left = Inches(0.8) + i * (p_w + Inches(0.3))
        c_top = Inches(1.7)
        add_card(s3, c_left, c_top, p_w, p_h)

        bar = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, c_left, c_top, p_w, Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = acc
        bar.line.fill.background()

        box = s3.shapes.add_textbox(c_left + Inches(0.2), c_top + Inches(0.25), p_w - Inches(0.4), p_h - Inches(0.5))
        btf = box.text_frame
        btf.word_wrap = True

        p0 = btf.paragraphs[0]
        p0.text = badge
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = acc
        p0.space_after = Pt(2)

        p1 = btf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE

        p1_sub = btf.add_paragraph()
        p1_sub.text = subtitle
        p1_sub.font.size = Pt(11)
        p1_sub.font.color.rgb = acc
        p1_sub.space_after = Pt(12)

        p2 = btf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 4: Pillar 1 — Knowledge Graph (Split Screen with Image)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    add_background(s4)
    add_header(s4, "Pillar 1: Heterogeneous Graph & Money Mule Ring Topologies")

    # Left: Explanation Points
    add_card(s4, Inches(0.8), Inches(1.7), Inches(6.4), Inches(5.1))
    box_s4 = s4.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(6.0), Inches(4.7))
    tf_s4 = box_s4.text_frame
    tf_s4.word_wrap = True

    g_points = [
        ("🔄 Tarjan's Directed Cycle Detection (Wash Trading)", "Detects circular money wash loops (A -> B -> C -> D -> A) in real-time O(V+E) time to uncover volume inflation and money laundering circuits.", CYAN),
        ("⚡ Rapid Drain / Pass-Through Mule Interceptor", "Mathematical signature: Outflow >= 85% of fresh inflow within Delta_T <= 15 minutes. Automatically flags intermediary transit mules.", AMBER),
        ("📱 Multi-Hop Hardware Syndicate Resolution", "Links disparate accounts sharing identical browser canvas hashes, IMEIs, or /24 subnet footprints into connected criminal syndicates.", PURPLE),
        ("💥 Topological Blast Radius Calculation", "Quantifies the total monetary contagion (₹ capital at risk) across 2-hop neighbor networks during syndicate outbreaks.", RED)
    ]

    for idx, (title, text, acc) in enumerate(g_points):
        p_t = tf_s4.paragraphs[0] if idx == 0 else tf_s4.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = acc

        p_b = tf_s4.add_paragraph()
        p_b.text = text
        p_b.font.size = Pt(11.5)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_after = Pt(8)

    # Right: Embedded Graphic
    if os.path.exists(IMG_GRAPH):
        add_card(s4, Inches(7.4), Inches(1.7), Inches(5.1), Inches(5.1))
        s4.shapes.add_picture(IMG_GRAPH, Inches(7.5), Inches(1.8), width=Inches(4.9), height=Inches(4.9))

    # =========================================================================
    # SLIDE 5: Pillar 2 — Behavioral Biometrics (Split Screen with Image)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    add_background(s5)
    add_header(s5, "Pillar 2: Behavioral Biometrics & Pre-Transaction Telemetry")

    # Left: Embedded Image
    if os.path.exists(IMG_BIO):
        add_card(s5, Inches(0.8), Inches(1.7), Inches(5.1), Inches(5.1))
        s5.shapes.add_picture(IMG_BIO, Inches(0.9), Inches(1.8), width=Inches(4.9), height=Inches(4.9))

    # Right: Explanation Points
    add_card(s5, Inches(6.1), Inches(1.7), Inches(6.4), Inches(5.1))
    box_s5 = s5.shapes.add_textbox(Inches(6.3), Inches(1.9), Inches(6.0), Inches(4.7))
    tf_s5 = box_s5.text_frame
    tf_s5.word_wrap = True

    bio_items = [
        ("⌨️ Keystroke Dynamics & Cadence Entropy", "Analyzes key flight times and dwell times. Automated bot scripts exhibit sub-human variance (<6ms) compared to natural human typing jitter.", EMERALD),
        ("📋 Real-Time Clipboard Paste Intercept (BioCatch)", "92% of APP scam victims copy-paste account numbers from scam phone/WhatsApp chats. Real-time DOM listeners trigger interactive scam warnings.", RED),
        ("✈️ Haversine Impossible Travel Velocity", "Calculates great-circle physical speed between consecutive GPS/IP coordinates. Physical speeds exceeding 800 km/h trigger instant ATO locks.", CYAN),
        ("⏱️ Hesitation & Form Dwelling Index", "Measures dwell pauses and cursor hesitation bursts prior to authorization, detecting active social engineering coercion.", AMBER)
    ]

    for idx, (title, text, acc) in enumerate(bio_items):
        p_t = tf_s5.paragraphs[0] if idx == 0 else tf_s5.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = acc

        p_b = tf_s5.add_paragraph()
        p_b.text = text
        p_b.font.size = Pt(11.5)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_after = Pt(8)

    # =========================================================================
    # SLIDE 6: Pillar 3 — Watchlist Screening & Explainable Risk Scoring
    # =========================================================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    add_background(s6)
    add_header(s6, "Pillar 3: Sanctions Screening & Explainable Multi-Vector Scoring")

    # Left: Formula Card
    add_card(s6, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.1))
    box_l = s6.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(5.2), Inches(4.7))
    tf_l = box_l.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "📐 Explainable Dynamic Risk Formula (0-100)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_after = Pt(12)

    weights = [
        ("Graph Topology Risk (30%)", "Cycles, Rapid Drain, In-Degree Centrality"),
        ("Behavioral Biometrics (25%)", "Keystroke Entropy, Paste Detection, Geo-Velocity"),
        ("Sanctions & PEP (20%)", "Jaro-Winkler Watchlist Hit Confidence"),
        ("AML Structuring (15%)", "Sub-Threshold PAN Limit Clustering (e.g. ₹49,999)"),
        ("Balance Depletion (10%)", ">90% Account Balance Drain Ratio")
    ]
    for w_title, w_desc in weights:
        wp = tf_l.add_paragraph()
        wp.text = f"• {w_title}: {w_desc}"
        wp.font.size = Pt(12)
        wp.font.color.rgb = TEXT_MUTED
        wp.space_after = Pt(6)

    # Right: Reason Codes Card
    add_card(s6, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1))
    box_r = s6.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.3), Inches(4.7))
    tf_r = box_r.text_frame
    tf_r.word_wrap = True

    pr0 = tf_r.paragraphs[0]
    pr0.text = "⚖️ Jaro-Winkler Sanctions & Reason Codes"
    pr0.font.size = Pt(16)
    pr0.font.bold = True
    pr0.font.color.rgb = PURPLE
    pr0.space_after = Pt(12)

    codes = [
        ("OFAC & PEP Fuzzy Matching", "Jaro-Winkler phonetic similarity matches obfuscated aliases (e.g. 'Al-Hassan Corp' vs 'Alhasan Trade')."),
        ("Standardized Reason Codes", "MULE_RAPID_DRAIN, DEVICE_COLLUSION_SYNDICATE, PASTE_ATTACK_APP_FRAUD, SANCTIONS_WATCHLIST_HIT."),
        ("Dynamic Tier Actions", "Low (<35) -> Instant Allow\nMed (35-74) -> WhatsApp OTP Challenge\nHigh (>=75) -> Auto-Quarantine & SAR Filing.")
    ]
    for c_title, c_desc in codes:
        cp = tf_r.add_paragraph()
        cp.text = f"• {c_title}:\n  {c_desc}"
        cp.font.size = Pt(12)
        cp.font.color.rgb = TEXT_MUTED
        cp.space_after = Pt(8)

    # =========================================================================
    # SLIDE 7: Pillar 4 — Autonomous FIU War Room (Split Screen with Image)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_slide_layout)
    add_background(s7)
    add_header(s7, "Pillar 4: Autonomous AI Financial Intelligence Unit & SAR Filing")

    # Left: Explanation Points
    add_card(s7, Inches(0.8), Inches(1.7), Inches(6.4), Inches(5.1))
    box_s7 = s7.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(6.0), Inches(4.7))
    tf_s7 = box_s7.text_frame
    tf_s7.word_wrap = True

    fiu_items = [
        ("🔍 FORENSIC AGENT (Topology Crawler)", "Crawls multi-hop graph relationships, compiles transaction chronologies, correlates hardware footprints, and computes infected blast radius.", CYAN),
        ("☣️ CONTAINMENT AGENT (Blast-Radius Quarantine)", "Executes 1-click simultaneous liquidity freezes across all correlated accounts in the contagion ring to prevent fund flight.", RED),
        ("📄 COMPLIANCE AGENT (Automated SAR/STR Filing)", "Auto-drafts official Suspicious Activity Reports (SAR / STR) with complete legal narratives and statutory citations (BSA 31 U.S.C. 5318(g), PMLA).", PURPLE)
    ]

    for idx, (title, text, acc) in enumerate(fiu_items):
        p_t = tf_s7.paragraphs[0] if idx == 0 else tf_s7.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = acc

        p_b = tf_s7.add_paragraph()
        p_b.text = text
        p_b.font.size = Pt(11.5)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_after = Pt(10)

    # Right: Embedded Command Center Image
    if os.path.exists(IMG_WAR_ROOM):
        add_card(s7, Inches(7.4), Inches(1.7), Inches(5.1), Inches(5.1))
        s7.shapes.add_picture(IMG_WAR_ROOM, Inches(7.5), Inches(1.8), width=Inches(4.9), height=Inches(4.9))

    # =========================================================================
    # SLIDE 8: Business Impact & Live Demo Highlights
    # =========================================================================
    s8 = prs.slides.add_slide(blank_slide_layout)
    add_background(s8)
    add_header(s8, "Business Impact, Benchmarks & Live Demonstration")

    # Metrics Banner
    stats = [
        ("< 15 ms", "Real-Time Decision Latency", CYAN),
        ("92% Drop", "Manual SAR Filing Effort", EMERALD),
        ("100% Contagion", "Blast Radius Containment", PURPLE),
        ("Zero Hardcoding", "Dynamic Graph Inference", BLUE)
    ]
    sw = Inches(2.7)
    sh = Inches(1.8)
    for i, (metric, label, acc) in enumerate(stats):
        c_left = Inches(0.8) + i * (sw + Inches(0.3))
        add_card(s8, c_left, Inches(1.7), sw, sh)

        box = s8.shapes.add_textbox(c_left + Inches(0.15), Inches(1.85), sw - Inches(0.3), sh - Inches(0.3))
        btf = box.text_frame
        btf.word_wrap = True

        p0 = btf.paragraphs[0]
        p0.text = metric
        p0.font.size = Pt(22)
        p0.font.bold = True
        p0.font.color.rgb = acc

        p1 = btf.add_paragraph()
        p1.text = label
        p1.font.size = Pt(11)
        p1.font.color.rgb = TEXT_MUTED

    # Demo Flow Steps
    add_card(s8, Inches(0.8), Inches(3.8), Inches(11.7), Inches(3.0))
    box_demo = s8.shapes.add_textbox(Inches(1.1), Inches(4.0), Inches(11.1), Inches(2.6))
    tf_d = box_demo.text_frame
    tf_d.word_wrap = True

    pd0 = tf_d.paragraphs[0]
    pd0.text = "🧪 Live Verification & Demonstration Flow:"
    pd0.font.size = Pt(16)
    pd0.font.bold = True
    pd0.font.color.rgb = TEXT_WHITE
    pd0.space_after = Pt(10)

    steps = [
        ("1. User Dashboard (/user/dashboard)", "Initiate transfer & paste recipient -> BioCatch scam alert triggers immediately."),
        ("2. Multi-Entity Intelligence Graph (/network)", "Inspect live heterogeneous graph with accounts, devices, IPs, and Tarjan wash cycles."),
        ("3. OFAC / PEP Sanctions Screener", "Fuzzy search 'Al-Hassan' or 'Viktor Antonov' to demonstrate Jaro-Winkler phonetic matching."),
        ("4. Autonomous FIU War Room", "Click 'Generate SAR' to view official FinCEN/RBI legal filing or trigger 'Quarantine Blast Radius'.")
    ]
    for s_title, s_desc in steps:
        sp = tf_d.add_paragraph()
        sp.text = f"• {s_title}: {s_desc}"
        sp.font.size = Pt(12)
        sp.font.color.rgb = TEXT_MUTED
        sp.space_after = Pt(4)

    # Save presentation
    output_filename = "Salaar_Bank_FinCrime_OS_Deck.pptx"
    try:
        prs.save(output_filename)
        print(f"[SUCCESS] Creative Presentation with visual assets generated: {os.path.abspath(output_filename)}")
    except Exception as e:
        alt_filename = f"Salaar_Bank_FinCrime_OS_Deck_{os.getpid()}.pptx"
        prs.save(alt_filename)
        print(f"[SUCCESS] Creative Presentation saved as fallback: {os.path.abspath(alt_filename)}")

if __name__ == '__main__':
    create_creative_deck()
